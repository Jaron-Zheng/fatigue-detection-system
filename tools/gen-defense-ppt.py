# -*- coding: utf-8 -*-
"""
gen-defense-ppt.py — 生成答辩演示 PPT（16:9，11 页）
素材复用 screenshots/round5-final/ 的实测截图，数字与 web/js/config.js、docs/测试报告.md 一致。
输出: docs/答辩演示.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOT = os.path.join(BASE, 'screenshots', 'round5-final')
OUT = os.path.join(BASE, 'docs', '答辩演示.pptx')

# ---------- 调色板（与项目令牌一致） ----------
INK = RGBColor(0x1A, 0x1D, 0x24)      # 主文字
SUB = RGBColor(0x5F, 0x66, 0x72)      # 次文字
FAINT = RGBColor(0x9A, 0xA1, 0xAD)    # 弱文字
BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF5, 0xF6, 0xF8)     # 浅灰卡
LINE = RGBColor(0xE4, 0xE7, 0xEB)
DARK = RGBColor(0x0B, 0x0F, 0x14)     # 封面深底
ACCENT = RGBColor(0x0A, 0x84, 0xFF)   # iOS 蓝
GREEN = RGBColor(0x30, 0xD1, 0x58)
YELLOW = RGBColor(0xFF, 0xD6, 0x0A)
ORANGE = RGBColor(0xFF, 0x9F, 0x0A)
RED = RGBColor(0xFF, 0x45, 0x3A)
LV = [GREEN, RGBColor(0xE3, 0xB3, 0x00), ORANGE, RED]  # 黄色加深便于白字
FONT = 'Microsoft YaHei'
MONO = 'Consolas'

EMU_IN = 914400

def C(hexstr):
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))

def _ea(run, name):
    """中文字体需要单独设置 a:ea，否则部分播放器回退宋体"""
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn('a:latin'))
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {'typeface': name})
        if latin is not None:
            latin.addnext(ea)
        else:
            rPr.append(ea)
    else:
        ea.set('typeface', name)

def style(run, size, bold=False, color=INK, name=FONT, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = name
    _ea(run, name)

def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: [ (text, size, bold, color, name, space_after_pt, line_spacing, align_override) ]"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for ln in lines:
        text, size, bold, color, name, sp_after, lsp, al = (ln + (None,) * 8)[:8]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = al if al is not None else align
        if sp_after is not None:
            p.space_after = Pt(sp_after)
        if lsp is not None:
            p.line_spacing = lsp
        r = p.add_run()
        r.text = text
        style(r, size, bold, color, name or FONT)
    return tb

def box(slide, x, y, w, h, fill=CARD, line=None, radius=0.09, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=0.75):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def chip(slide, x, y, w, h, text, fill, tcolor=RGBColor(0xFF, 0xFF, 0xFF), size=11, bold=True, line=None):
    sp = box(slide, x, y, w, h, fill=fill, line=line, radius=0.5)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    style(r, size, bold, tcolor)
    return sp

def header(slide, num, title, subtitle=None):
    box(slide, 0.6, 0.52, 0.09, 0.62, fill=ACCENT, shape=MSO_SHAPE.RECTANGLE)
    txt(slide, 0.82, 0.44, 10.5, 0.5,
        [(f'{num:02d}', 13, True, ACCENT, FONT, 0, 1.0)])
    txt(slide, 0.82, 0.72, 10.5, 0.62,
        [(title, 26, True, INK, FONT, 0, 1.0)])
    if subtitle:
        txt(slide, 0.82, 1.32, 11.9, 0.34,
            [(subtitle, 12.5, False, SUB, FONT, 0, 1.0)])
    txt(slide, 9.7, 0.5, 3.0, 0.3,
        [('疲劳检测系统 · 毕业设计答辩', 10, False, FAINT, FONT, 0, 1.0)],
        align=PP_ALIGN.RIGHT)

def footer(slide, n, total=11):
    txt(slide, 0.6, 7.08, 8.0, 0.3,
        [('基于面部多特征融合的 Web 端驾驶员疲劳检测系统', 9, False, FAINT, FONT, 0, 1.0)])
    txt(slide, 12.2, 7.08, 0.55, 0.3,
        [(f'{n} / {total}', 9, False, FAINT, FONT, 0, 1.0)], align=PP_ALIGN.RIGHT)

def accent_strip(slide, x, y, w=3.2, h=0.07):
    seg = w / 4
    for i, c in enumerate(LV):
        box(slide, x + i * seg, y, seg - 0.04, h, fill=c, shape=MSO_SHAPE.RECTANGLE)

def picture_framed(slide, path, x, y, w, h, border=ACCENT):
    box(slide, x - 0.05, y - 0.05, w + 0.1, h + 0.1, fill=DARK, line=border, radius=0.03, line_w=1.2)
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

# ================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def new_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    box(s, 0, 0, 13.333, 7.5, fill=bg, shape=MSO_SHAPE.RECTANGLE)
    return s

# ---------------- P1 封面 ----------------
s = new_slide(DARK)
accent_strip(s, 0.62, 0.62, 3.2)
txt(s, 0.62, 0.92, 6.0, 0.36, [('本科毕业设计 · 答辩汇报', 13, False, RGBColor(0x9A, 0xA6, 0xB8), FONT, 0, 1.0)])
txt(s, 0.6, 1.78, 6.4, 2.2, [
    ('基于面部多特征融合的', 34, True, RGBColor(0xF2, 0xF4, 0xF8), FONT, 6, 1.12),
    ('Web 端驾驶员疲劳检测系统', 34, True, RGBColor(0xFF, 0xFF, 0xFF), FONT, 0, 1.12),
])
txt(s, 0.62, 3.62, 6.3, 0.4, [
    ('纯浏览器运行  ·  全本地计算  ·  零安装  ·  PWA 离线可用', 14, False, ACCENT, FONT, 0, 1.0)])
txt(s, 0.62, 4.78, 6.3, 1.6, [
    ('汇报人：某某某　　学号：XXXXXXXX', 13, False, RGBColor(0xC8, 0xCF, 0xD9), FONT, 8, 1.2),
    ('指导教师：某某某　教授', 13, False, RGBColor(0xC8, 0xCF, 0xD9), FONT, 8, 1.2),
    ('2026 年 X 月', 13, False, RGBColor(0xC8, 0xCF, 0xD9), FONT, 0, 1.2),
])
picture_framed(s, os.path.join(SHOT, 'home-1366x768.png'), 7.05, 2.05, 5.75, 3.23,
               border=RGBColor(0x2A, 0x33, 0x40))
txt(s, 7.05, 5.4, 5.75, 0.3, [('系统首页（实测截图）：Tesla HUD 风格驾驶舱预览与疲劳仪表', 10, False, RGBColor(0x8A, 0x94, 0xA3), FONT, 0, 1.0)])

# ---------------- P2 背景与问题 ----------------
s = new_slide()
header(s, 2, '研究背景与问题', '驾驶疲劳监测的三类现有方案，各有难以落地的痛点')
pains = [
    ('专用硬件方案', '方向盘握力 / 眼动仪 / 可穿戴设备', '装机成本高，后装渗透率低，难以大范围推广'),
    ('单一特征方案', '仅依赖 PERCLOS 等单一指标', '个体差异大，眼镜、姿态干扰下误报率高，判定不可解释'),
    ('云端视觉方案', '视频上传服务器进行分析', '存在隐私顾虑；弱网 / 断网环境完全不可用'),
]
y = 1.62
for i, (t, m, d) in enumerate(pains):
    box(s, 0.6, y, 7.55, 1.52, fill=CARD, line=LINE)
    box(s, 0.6, y + 0.18, 0.07, 1.16, fill=RED, shape=MSO_SHAPE.RECTANGLE)
    txt(s, 0.88, y + 0.16, 7.1, 0.4, [(t, 16, True, INK, FONT, 2, 1.0)])
    txt(s, 0.88, y + 0.56, 7.1, 0.32, [(m, 11.5, False, SUB, FONT, 0, 1.0)])
    txt(s, 0.88, y + 0.9, 7.1, 0.5, [(d, 12, False, SUB, FONT, 0, 1.08)])
    y += 1.68
box(s, 8.4, 1.62, 4.33, 3.1, fill=DARK)
txt(s, 8.72, 1.92, 3.7, 0.4, [('选题思路', 15, True, RGBColor(0xFF, 0xFF, 0xFF), FONT, 0, 1.0)])
accent_strip(s, 8.72, 2.34, 1.6)
ideas = ['多特征融合决策——不押注单一指标', '全本地计算——视频数据不出本机', '零安装——打开浏览器即可使用']
yy = 2.56
for it in ideas:
    txt(s, 8.72, yy, 0.3, 0.3, [('▶', 9, True, ACCENT, FONT, 0, 1.0)])
    txt(s, 9.0, yy, 3.55, 0.62, [(it, 11.5, False, RGBColor(0xD5, 0xDB, 0xE3), FONT, 0, 1.12)])
    yy += 0.66
box(s, 8.4, 4.92, 4.33, 1.62, fill=RGBColor(0xEA, 0xF2, 0xFF), line=RGBColor(0xBF, 0xD9, 0xFF))
txt(s, 8.72, 5.12, 3.75, 1.3, [
    ('核心问题', 12, True, ACCENT, FONT, 4, 1.0),
    ('能否不装任何硬件、不上传任何视频，在浏览器里做出可信、可解释的疲劳判定？', 12.5, True, INK, FONT, 0, 1.22),
])
footer(s, 2)

# ---------------- P3 目标与指标 ----------------
s = new_slide()
header(s, 3, '研究目标与技术指标', '输出 [0,100] 连续疲劳指数，并划分为四级疲劳状态')
lx, lw = 0.6, 12.13
segs = [(30, '清醒', LV[0], '0 – 30'), (22, '轻度疲劳', LV[1], '30 – 52'), (22, '中度疲劳', LV[2], '52 – 74'), (26, '重度疲劳', LV[3], '74 – 100')]
cx = lx
txt(s, lx, 1.58, 6.0, 0.3, [('疲劳指数分级（阈值与 web/js/config.js 完全一致）', 11, False, SUB, FONT, 0, 1.0)])
for w_, label, c, rng in segs:
    ww = lw * w_ / 100
    sp = box(s, cx, 1.92, ww - 0.05, 0.78, fill=c)
    tf = sp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    style(r, 15, True, RGBColor(0xFF, 0xFF, 0xFF) if w_ != 30 else INK)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = rng
    style(r2, 10, False, RGBColor(0xFF, 0xFF, 0xFF) if w_ != 30 else INK)
    cx += ww
kpis = [
    ('实时性', '单帧推理 ≤ 45ms', 'GPU / CPU 双委托，GPU 失败自动回退（e2e 实测）'),
    ('隐私性', '数据全程不出本机', '模型与 WASM 全本地加载，无第三方运行时请求'),
    ('可用性', 'PWA 离线全功能', '断网下检测、报警、导出均实测可用（12/12 断网用例）'),
    ('兼容性', '三大引擎实测通过', 'Chromium / Firefox / WebKit 各 12 项断言，24/24'),
]
x = 0.6
for t, v, d in kpis:
    box(s, x, 3.3, 2.92, 2.5, fill=CARD, line=LINE)
    txt(s, x + 0.24, 3.56, 2.5, 0.36, [(t, 12, True, ACCENT, FONT, 0, 1.0)])
    txt(s, x + 0.24, 3.98, 2.5, 0.8, [(v, 16, True, INK, FONT, 0, 1.1)])
    txt(s, x + 0.24, 4.78, 2.5, 0.95, [(d, 10.5, False, SUB, FONT, 0, 1.15)])
    x += 3.07
box(s, 0.6, 6.1, 12.13, 0.62, fill=RGBColor(0xEA, 0xF2, 0xFF), line=RGBColor(0xBF, 0xD9, 0xFF))
txt(s, 0.88, 6.24, 11.6, 0.36, [
    ('判定语义：轻度=提示提醒 · 中度=建议休息 · 重度=强报警（声光+横幅，分级冷却防轰炸）', 12, True, INK, FONT, 0, 1.0)])
footer(s, 3)

# ---------------- P4 系统架构 ----------------
s = new_slide()
header(s, 4, '系统架构：三层流水线 + 会话状态机', '纯前端单页应用；MediaPipe WASM 推理；状态机守护全部交互')
layers = [
    ('采集层', ACCENT, ['getUserMedia 30fps 视频流', 'MediaPipe FaceLandmarker', '478 个人脸关键点', 'WebGL / CPU 双委托推理']),
    ('特征层', ORANGE, ['EAR 眼纵横比 · MAR 嘴纵横比', 'PERCLOS 时间加权占比', '欧拉角头部姿态解算', '眨眼 / 哈欠 / 点头 / 分神事件']),
    ('决策层', RED, ['模糊加权融合引擎', '迟滞 + 驻留防抖判定', '语义否决假阳性', '四级预警 · 分级声光报警']),
]
x = 0.6
for t, c, items in layers:
    box(s, x, 1.72, 3.72, 3.0, fill=CARD, line=LINE)
    chip(s, x + 0.24, 1.94, 1.35, 0.42, t, c, size=13)
    yy = 2.56
    for it in items:
        txt(s, x + 0.28, yy, 3.2, 0.4, [('· ', 11.5, True, c, FONT, 0, 1.0)])
        txt(s, x + 0.5, yy, 3.05, 0.44, [(it, 11.5, False, INK, FONT, 0, 1.05)])
        yy += 0.5
    if x < 8.5:
        ar = box(s, x + 3.78, 2.96, 0.36, 0.5, fill=RGBColor(0xC5, 0xCD, 0xD8), shape=MSO_SHAPE.CHEVRON)
    x += 4.2
box(s, 0.6, 5.06, 12.13, 1.62, fill=DARK)
txt(s, 0.9, 5.24, 4.5, 0.36, [('会话状态机（守护一切交互的前提）', 12.5, True, RGBColor(0xFF, 0xFF, 0xFF), FONT, 0, 1.0)])
states = ['IDLE', 'BOOTING', 'CALIBRATING', 'RUNNING', 'PAUSED', 'REPORT']
sx = 0.9
for i, st in enumerate(states):
    c = ACCENT if st == 'RUNNING' else RGBColor(0x2A, 0x33, 0x40)
    chip(s, sx, 5.72, 1.52, 0.44, st, c, size=10.5, line=RGBColor(0x3A, 0x44, 0x52))
    if i < len(states) - 1:
        txt(s, sx + 1.54, 5.76, 0.42, 0.36, [('→', 13, True, RGBColor(0x8A, 0x94, 0xA3), FONT, 0, 1.0)])
    sx += 1.96
txt(s, 8.2, 5.28, 4.3, 1.2, [
    ('状态机 38 项断言全覆盖：', 11, True, RGBColor(0xD5, 0xDB, 0xE3), FONT, 3, 1.15),
    ('非法迁移一律拒绝并保持原状态；', 10.5, False, RGBColor(0x9A, 0xA4, 0xB2), FONT, 0, 1.15),
    ('守卫缺失类 bug（未完成可导出等）由此根治', 10.5, False, RGBColor(0x9A, 0xA4, 0xB2), FONT, 0, 1.15),
])
footer(s, 4)

# ---------------- P5 特征提取 ----------------
s = new_slide()
header(s, 5, '算法核心一：多特征提取', '眼部 / 口部 / 姿态 / 事件四路信号，全部由 478 关键点实时导出')
feats = [
    ('EAR 眼纵横比', 'EAR = (|p2-p6| + |p3-p5|) / (2|p1-p4|)', '闭眼判定 = 睁眼基线 × 0.72；恢复 = × 0.80，双阈值滞回抗抖'),
    ('MAR 嘴纵横比 → 哈欠', '张口度持续超阈值且满足最短时长', '计一次哈欠事件，带不应期约束，过滤说话 / 深呼吸'),
    ('PERCLOS 闭眼时间占比', 'P80 标准；滑动窗口按区间真实时长加权', '采样率波动不引入偏置；间断观测不假设状态延续'),
    ('头部姿态（欧拉角）', '旋转矩阵 → yaw / pitch / roll', '点头 = 俯仰角速度峰值；分神 = 偏离持续超时'),
]
pos = [(0.6, 1.66), (6.75, 1.66), (0.6, 3.42), (6.75, 3.42)]
for (t, f, d), (px, py) in zip(feats, pos):
    box(s, px, py, 5.98, 1.62, fill=CARD, line=LINE)
    txt(s, px + 0.26, py + 0.16, 5.5, 0.38, [(t, 14.5, True, INK, FONT, 0, 1.0)])
    txt(s, px + 0.26, py + 0.56, 5.55, 0.4, [(f, 11, True, ACCENT, MONO, 0, 1.05)])
    txt(s, px + 0.26, py + 1.0, 5.55, 0.55, [(d, 10.5, False, SUB, FONT, 0, 1.12)])
box(s, 0.6, 5.24, 12.13, 1.32, fill=RGBColor(0xE8, 0xF8, 0xEE), line=RGBColor(0xB7, 0xE6, 0xC6))
txt(s, 0.9, 5.42, 3.2, 0.4, [('个性化校准（8 秒）', 14, True, RGBColor(0x1D, 0x8A, 0x43), FONT, 0, 1.0)])
txt(s, 0.9, 5.86, 11.6, 0.62, [
    ('检测前采集 ≥ 60 个有效样本建立睁眼基线，闭眼阈值随个体自适应——从机制上消解“眼睛大小 / 眼睑形态”个体差异，这是传统固定阈值方案误报高的根源。', 11.5, False, INK, FONT, 0, 1.15)])
footer(s, 5)

# ---------------- P6 融合引擎 ----------------
s = new_slide()
header(s, 6, '算法核心二：模糊加权融合引擎', '七特征隶属度加权求和 → 连续指数 → 迟滞分级 → 分级报警')
box(s, 0.6, 1.66, 6.1, 3.1, fill=CARD, line=LINE)
txt(s, 0.88, 1.84, 5.6, 0.42, [('疲劳指数 = 100 × Σ wk·μk(xk)', 16, True, INK, FONT, 0, 1.0)])
txt(s, 0.88, 2.3, 5.6, 0.3, [('μk 为各特征对“疲劳”的隶属度，wk 权重合计 = 1.00（config.js）', 10, False, SUB, FONT, 0, 1.0)])
weights = [
    ('PERCLOS 闭眼占比', 0.30, RED), ('最长持续闭眼', 0.20, ORANGE), ('哈欠频率', 0.14, YELLOW),
    ('眨眼频率偏移', 0.10, ACCENT), ('点头频率', 0.10, ACCENT), ('平均眨眼时长', 0.08, GREEN), ('头部偏离占比', 0.08, GREEN),
]
yy = 2.68
for name, w_, c in weights:
    txt(s, 0.88, yy, 2.4, 0.3, [(name, 10.5, False, INK, FONT, 0, 1.0)])
    box(s, 3.35, yy + 0.055, 2.2 * w_ / 0.30, 0.2, fill=c, shape=MSO_SHAPE.RECTANGLE)
    txt(s, 5.75, yy, 0.8, 0.3, [(f'{w_:.2f}', 10.5, True, INK, MONO, 0, 1.0)])
    yy += 0.29
mechs = [
    ('迟滞 + 驻留', '跌回下一级需低 6 分，且需连续满足驻留时长——等级不抖动', RED),
    ('语义否决', '头部姿态异常时否决 EAR 假阳性（低头闭眼≠入睡），回归实证 0 误报', ORANGE),
    ('EMA 平滑', '指数平滑抑制瞬时毛刺，输出稳定可解释的连续曲线', ACCENT),
]
yy = 1.66
for t, d, c in mechs:
    box(s, 6.95, yy, 5.78, 0.98, fill=CARD, line=LINE)
    box(s, 6.95, yy + 0.14, 0.07, 0.7, fill=c, shape=MSO_SHAPE.RECTANGLE)
    txt(s, 7.2, yy + 0.12, 5.3, 0.34, [(t, 13, True, INK, FONT, 0, 1.0)])
    txt(s, 7.2, yy + 0.48, 5.4, 0.46, [(d, 10.5, False, SUB, FONT, 0, 1.1)])
    yy += 1.08
box(s, 6.95, 4.96, 5.78, 1.7, fill=DARK)
txt(s, 7.25, 5.14, 5.2, 0.36, [('决策链路', 12.5, True, RGBColor(0xFF, 0xFF, 0xFF), FONT, 0, 1.0)])
chain = ['连续指数（EMA）', '等级判定（迟滞+驻留）', '分级报警（冷却）']
sx = 7.25
for i, st in enumerate(chain):
    chip(s, sx, 5.62, 1.72, 0.5, st, RGBColor(0x1E, 0x28, 0x36), size=9.5, line=ACCENT)
    if i < 2:
        txt(s, sx + 1.74, 5.68, 0.3, 0.36, [('→', 12, True, ACCENT, FONT, 0, 1.0)])
    sx += 1.86
txt(s, 7.25, 6.22, 5.2, 0.34, [('全部参数有文献依据（见《算法参数文献核查记录》）', 9.5, False, RGBColor(0x9A, 0xA4, 0xB2), FONT, 0, 1.0)])
box(s, 0.6, 4.96, 6.1, 1.7, fill=RGBColor(0xEA, 0xF2, 0xFF), line=RGBColor(0xBF, 0xD9, 0xFF))
txt(s, 0.88, 5.14, 5.55, 0.36, [('为什么用模糊逻辑而不是深度学习？', 12.5, True, ACCENT, FONT, 0, 1.0)])
txt(s, 0.88, 5.56, 5.6, 1.0, [
    ('无需大规模标注训练数据；权重与阈值可解释、可审计、可个性化调整——', 10.5, False, INK, FONT, 0, 1.18),
    ('契合车载安全场景对“判定依据可追溯”的硬要求', 10.5, True, INK, FONT, 0, 1.18)])
footer(s, 6)

# ---------------- P7 工程亮点 ----------------
s = new_slide()
header(s, 7, '工程实现亮点', '把“演示 Demo”做成“可交付系统”的四项关键工程')
hl = [
    ('推理引擎韧性', 'WebGL GPU 委托失败自动回退 CPU；假摄像头 e2e 双链路各 11 项断言实测通过，授权 15s 超时兜底'),
    ('零外部依赖', '模型 / WASM 全本地；CSP default-src \'self\'；首页无任何外部域名请求（离线可用前提）'),
    ('配置安全三重校验', 'localStorage 载入配置经原型污染拦截 → 类型校验 → 47 条路径数值区间钳制（越界收边界）'),
    ('工程化质量门禁', 'verify:full 一键全量：静态检查 + 回归 137 + 集成 41 + typecheck + lint，全绿才算过'),
]
pos = [(0.6, 1.72), (6.75, 1.72), (0.6, 3.72), (6.75, 3.72)]
for (t, d), (px, py) in zip(hl, pos):
    box(s, px, py, 5.98, 1.8, fill=CARD, line=LINE)
    chip(s, px + 0.26, py + 0.2, 0.5, 0.5, '✓', ACCENT, size=15)
    txt(s, px + 0.92, py + 0.22, 4.9, 0.42, [(t, 15, True, INK, FONT, 0, 1.0)])
    txt(s, px + 0.92, py + 0.68, 4.85, 1.0, [(d, 11, False, SUB, FONT, 0, 1.18)])
box(s, 0.6, 5.66, 12.13, 1.24, fill=DARK)
txt(s, 0.9, 5.82, 4.0, 0.36, [('安全纵深（集成测试 41 项实证）', 12.5, True, RGBColor(0xFF, 0xFF, 0xFF), FONT, 0, 1.0)])
sec = [
    '目录穿越 / URL 编码穿越 / 空字节注入 → 全部拒绝',
    'CSV 公式注入防护；XSS 全输入面测试通过',
    'COOP / Permissions-Policy 收紧摄像头权限面',
]
sxx = 5.1
for it in sec:
    txt(s, sxx, 5.8, 0.25, 0.3, [('·', 12, True, ACCENT, FONT, 0, 1.0)])
    txt(s, sxx + 0.22, 5.82, 2.5, 0.92, [(it, 10, False, RGBColor(0xD5, 0xDB, 0xE3), FONT, 0, 1.15)])
    sxx += 2.62
footer(s, 7)

# ---------------- P8 质量保障 ----------------
s = new_slide()
header(s, 8, '质量保障体系：对抗性测试', '方法论：交互清单普查 × 状态组合矩阵 × 反逻辑用户角色扮演')
stats = [('400+', '项自动化断言'), ('20+', '套可复跑脚本'), ('0', '控制台错误'), ('0', 'axe 严重问题')]
x = 0.6
for v, t in stats:
    box(s, x, 1.6, 2.92, 1.02, fill=DARK)
    txt(s, x, 1.72, 2.92, 0.52, [(v, 26, True, RGBColor(0xFF, 0xFF, 0xFF), FONT, 0, 1.0)], align=PP_ALIGN.CENTER)
    txt(s, x, 2.24, 2.92, 0.3, [(t, 10.5, False, RGBColor(0x9A, 0xA4, 0xB2), FONT, 0, 1.0)], align=PP_ALIGN.CENTER)
    x += 3.07
suites = [
    ('单元回归', 137, GREEN), ('集成/安全', 41, GREEN), ('UI 全流程', 56, GREEN), ('交互混沌', 20, GREEN),
    ('开关混沌', 21, GREEN), ('跨引擎', 24, GREEN), ('状态边界', 22, GREEN), ('导出一致性', 22, GREEN),
    ('纯键盘', 9, GREEN), ('中断恢复', 7, GREEN), ('故障注入', 6, GREEN), ('探索角色', 19, GREEN),
]
txt(s, 0.6, 2.78, 12.0, 0.32, [('测试套件全景（全部通过，tools/ 下可一键复跑）', 11.5, True, INK, FONT, 0, 1.0)])
x, y = 0.6, 3.16
for i, (name, n, c) in enumerate(suites):
    box(s, x, y, 2.92, 0.6, fill=CARD, line=LINE)
    txt(s, x + 0.18, y + 0.14, 2.0, 0.34, [(name, 11, True, INK, FONT, 0, 1.0)])
    chip(s, x + 2.14, y + 0.11, 0.62, 0.38, str(n), c, size=11)
    x += 3.07
    if i % 4 == 3:
        x = 0.6
        y += 0.7
box(s, 0.6, 5.32, 12.13, 1.5, fill=RGBColor(0xFF, 0xF3, 0xF0), line=RGBColor(0xFF, 0xC7, 0xBF))
txt(s, 0.9, 5.46, 11.6, 0.36, [
    ('典型案例（现象 A · P0）：专业 / 普通模式下载的报告内容相同', 13, True, RED, FONT, 0, 1.0)])
txt(s, 0.9, 5.88, 11.7, 0.9, [
    ('根因：HTML 导出的“空壳折叠”逻辑不区分模式——普通版泄漏专业术语，专业版折成同样文案，两份文件肉眼一致。', 11, False, INK, FONT, 3, 1.18),
    ('演进：先按模式分流修复；后按交付需求升级为统一专业版——无论开关，导出恒为详细报告（body 强制 pro-mode），CSV / JSON 与模式无关。开关混沌 24/24 复测通过。', 11, False, INK, FONT, 0, 1.18)])
footer(s, 8)

# ---------------- P9 现场演示 ----------------
s = new_slide()
header(s, 9, '现场演示（演示模式，60 秒）', '四级预警全流程：与真实摄像头同链路，仅信号源换为模拟剧本')
steps = [
    ('1', '启动检测', '浏览器打开即用，无任何安装'),
    ('2', '校准 8 秒', '建立个体睁眼基线'),
    ('3', '实时检测', '示意脸 + 指标曲线 + 特征值联动'),
    ('4', '四级预警', '轻度提示 → 重度声光报警横幅'),
    ('5', '报告导出', 'HTML / JSON / CSV 三格式秒级生成'),
]
y = 1.72
for n, t, d in steps:
    chip(s, 0.6, y + 0.06, 0.44, 0.44, n, ACCENT, size=13)
    txt(s, 1.2, y + 0.02, 3.1, 0.36, [(t, 13.5, True, INK, FONT, 0, 1.0)])
    txt(s, 1.2, y + 0.4, 3.3, 0.36, [(d, 10.5, False, SUB, FONT, 0, 1.0)])
    if n != '5':
        box(s, 0.79, y + 0.56, 0.06, 0.5, fill=LINE, shape=MSO_SHAPE.RECTANGLE)
    y += 1.02
picture_framed(s, os.path.join(SHOT, 'workbench-light.png'), 4.75, 1.7, 7.95, 4.47)
txt(s, 4.75, 6.24, 7.95, 0.3, [('工作台实测截图：实时示意脸 · 疲劳指数仪表 · 四特征曲线联动', 10, False, FAINT, FONT, 0, 1.0)])
box(s, 0.6, 6.56, 3.9, 0.5, fill=RGBColor(0xE8, 0xF8, 0xEE), line=RGBColor(0xB7, 0xE6, 0xC6))
txt(s, 0.82, 6.68, 3.6, 0.3, [('应急预案：断网照样演示（PWA 离线）', 10.5, True, RGBColor(0x1D, 0x8A, 0x43), FONT, 0, 1.0)])
footer(s, 9)

# ---------------- P10 局限与展望 ----------------
s = new_slide()
header(s, 10, '局限性与展望', '主动亮牌：如实呈现适用边界，比回避更专业')
txt(s, 0.6, 1.66, 5.9, 0.38, [('当前局限', 15, True, RED, FONT, 0, 1.0)])
lims = [
    ('光照与遮挡敏感', '逆光、墨镜、大面积口罩会显著降低关键点质量（质量门控会给出明确提示而非误判）'),
    ('纯视觉单模态', '无方向盘握力 / 心率等生理信号，模态间无法交叉印证'),
    ('验证规模有限', '自动化 + 假摄像头 e2e 充分，但真实驾驶舱大规模数据尚未采集'),
]
y = 2.12
for t, d in lims:
    box(s, 0.6, y, 5.9, 1.34, fill=CARD, line=LINE)
    txt(s, 0.88, y + 0.14, 5.4, 0.36, [(t, 13, True, INK, FONT, 0, 1.0)])
    txt(s, 0.88, y + 0.52, 5.35, 0.72, [(d, 10.5, False, SUB, FONT, 0, 1.15)])
    y += 1.48
txt(s, 6.85, 1.66, 5.9, 0.38, [('后续展望', 15, True, RGBColor(0x1D, 0x8A, 0x43), FONT, 0, 1.0)])
futs = [
    ('多模态融合', '接入方向盘握力 / 可穿戴心率，与视觉特征在同一模糊框架下加权'),
    ('车队端侧部署', 'PWA 形态天然适配车载 WebView，离线能力即车载数据不出车'),
    ('个性化在线学习', '把个体校准从“8 秒基线”升级为随使用持续更新的自适应基线'),
]
y = 2.12
for t, d in futs:
    box(s, 6.85, y, 5.9, 1.34, fill=RGBColor(0xE8, 0xF8, 0xEE), line=RGBColor(0xB7, 0xE6, 0xC6))
    txt(s, 7.13, y + 0.14, 5.4, 0.36, [(t, 13, True, INK, FONT, 0, 1.0)])
    txt(s, 7.13, y + 0.52, 5.35, 0.72, [(d, 10.5, False, SUB, FONT, 0, 1.15)])
    y += 1.48
footer(s, 10)

# ---------------- P11 致谢 ----------------
s = new_slide(DARK)
accent_strip(s, 5.07, 2.3, 3.2)
txt(s, 0.6, 2.72, 12.13, 0.9, [('感谢聆听，请各位老师批评指正', 34, True, RGBColor(0xFF, 0xFF, 0xFF), FONT, 0, 1.1)], align=PP_ALIGN.CENTER)
txt(s, 0.6, 3.78, 12.13, 0.4, [('Q & A', 18, True, ACCENT, FONT, 0, 1.0)], align=PP_ALIGN.CENTER)
recap = ['多特征模糊融合', '个体化校准', '400+ 断言质量门禁', 'PWA 离线可用']
x = 2.87
for it in recap:
    chip(s, x, 4.66, 1.95, 0.5, it, RGBColor(0x1E, 0x28, 0x36), size=10.5, line=RGBColor(0x3A, 0x44, 0x52))
    x += 2.0
txt(s, 0.6, 6.7, 12.13, 0.3, [('基于面部多特征融合的 Web 端驾驶员疲劳检测系统 · 2026', 10, False, RGBColor(0x8A, 0x94, 0xA3), FONT, 0, 1.0)], align=PP_ALIGN.CENTER)

# 画布边界自检：任何形状超出 13.333x7.5 即警告
W, H = 13.333, 7.5
issues = []
for idx, slide in enumerate(prs.slides, 1):
    for sh in slide.shapes:
        if sh.left is None:
            continue
        l, t = sh.left / EMU_IN, sh.top / EMU_IN
        r, b = l + sh.width / EMU_IN, t + sh.height / EMU_IN
        if r > W + 0.01 or b > H + 0.01 or l < -0.01 or t < -0.01:
            issues.append(f'slide{idx}: ({l:.2f},{t:.2f})-({r:.2f},{b:.2f})')
if issues:
    print('BOUNDS WARN:')
    for i in issues:
        print(' ', i)
else:
    print('bounds check: all shapes within canvas')

prs.save(OUT)
print(f'saved: {OUT}')
